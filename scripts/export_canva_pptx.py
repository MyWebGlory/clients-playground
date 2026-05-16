"""
Export static project slides to a Canva-compatible PPTX with movable layers.

How it works:
1. Render each slide in Playwright.
2. Capture a background layer (slide with direct children hidden).
3. Capture each direct child element as a transparent PNG layer.
4. Assemble one PPTX slide per HTML slide, adding each layer at matching coordinates.

Result:
- Canva can import the PPTX.
- Elements are movable (image layers).
- Text is flattened into image objects (no editable text boxes).

Usage:
  python scripts/export_canva_pptx.py --project conference-video-slides --selector .slide
  python scripts/export_canva_pptx.py --client cbhn --project conference-video-slides --selector .slide
    python scripts/export_canva_pptx.py --project conference-video-slides --selector .slide --slide-indices 1,2,3,4
"""

import argparse
import asyncio
import http.server
import os
import pathlib
import re
import shutil
import threading
import time
from dataclasses import dataclass
from typing import List, Optional, Set

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = pathlib.Path(__file__).resolve().parent.parent
PUBLIC = ROOT / "public"
TMP_ROOT = ROOT / ".tmp" / "canva-pptx-export"


@dataclass
class LayerCapture:
    image_path: pathlib.Path
    x: float
    y: float
    width: float
    height: float


@dataclass
class SlideCapture:
    source_index: int
    title: str
    width: float
    height: float
    layers: List[LayerCapture]


def project_exists(client_slug: str, project_slug: str) -> bool:
    return (PUBLIC / "clients" / client_slug / "projects" / project_slug / "index.html").exists()


def sanitize_filename(raw: str, extension: str) -> str:
    cleaned = re.sub(r"[\\/:*?\"<>|]", "-", raw or "")
    cleaned = re.sub(r"\s+", " ", cleaned).strip().strip(".")
    if not cleaned:
        cleaned = "Export"

    ext = extension if extension.startswith(".") else f".{extension}"
    if not cleaned.lower().endswith(ext.lower()):
        cleaned += ext

    return cleaned


def start_server(port: int):
    handler = http.server.SimpleHTTPRequestHandler
    handler.log_message = lambda *args: None
    server = http.server.HTTPServer(("127.0.0.1", port), handler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    return server


async def wait_for_assets(page):
    await page.evaluate("document.fonts && document.fonts.ready ? document.fonts.ready : Promise.resolve()")
    await page.evaluate(
        """
        Promise.all(Array.from(document.images).map((img) => {
            if (img.complete && img.naturalWidth > 0) {
                return img.decode ? img.decode().catch(() => undefined) : Promise.resolve();
            }
            return new Promise((resolve) => {
                const done = () => resolve();
                img.addEventListener('load', done, { once: true });
                img.addEventListener('error', done, { once: true });
            });
        }))
        """
    )
    await page.evaluate(
        """
        new Promise((resolve) => {
            requestAnimationFrame(() => requestAnimationFrame(resolve));
        })
        """
    )


async def detect_default_selector(page) -> Optional[str]:
    candidates = [".slide", ".graphic", ".flyer-page", ".page", ".zoom-graphic"]

    for selector in candidates:
        count = await page.locator(selector).count()
        if count > 0:
            return selector

    return None


async def extract_slide_title(slide) -> str:
    raw = await slide.evaluate(
        """
        (node) => {
            let previous = node.previousElementSibling;
            while (previous) {
                const label = previous.querySelector && previous.querySelector('.slide-label, .graphic-label, .section-label');
                const text = (label && label.textContent ? label.textContent : '').trim();
                if (text) {
                    return text;
                }
                previous = previous.previousElementSibling;
            }

            const heading = node.querySelector('h1, h2, h3, .title, .event-title');
            const headingText = (heading && heading.textContent ? heading.textContent : '').trim();
            if (headingText) {
                return headingText;
            }

            return '';
        }
        """
    )
    return str(raw or "").strip()


async def hide_direct_children(slide):
    await slide.evaluate(
        """
        (node) => {
            for (const child of Array.from(node.children)) {
                child.dataset.exportPrevVisibility = child.style.visibility || '';
                child.style.visibility = 'hidden';
            }
        }
        """
    )


async def restore_direct_children(slide):
    await slide.evaluate(
        """
        (node) => {
            for (const child of Array.from(node.children)) {
                child.style.visibility = child.dataset.exportPrevVisibility || '';
                delete child.dataset.exportPrevVisibility;
            }
        }
        """
    )


async def clear_capture_targets(slide):
    await slide.evaluate(
        """
        (node) => {
            for (const el of node.querySelectorAll('[data-export-target-id]')) {
                el.removeAttribute('data-export-target-id');
            }
        }
        """
    )


async def set_target_extra_bottom_padding(slide, target_id: str, extra_bottom_px: float, enabled: bool):
    await slide.evaluate(
        """
        (node, payload) => {
            const { targetId, extraBottomPx, enabled } = payload;
            const el = node.querySelector(`[data-export-target-id="${targetId}"]`);
            if (!el) {
                return;
            }

            const flag = 'data-export-extra-bottom-active';
            if (enabled) {
                if (el.hasAttribute(flag)) {
                    return;
                }

                const style = window.getComputedStyle(el);
                const basePaddingBottom = parseFloat(style.paddingBottom || '0') || 0;
                const baseMarginBottom = parseFloat(style.marginBottom || '0') || 0;

                el.dataset.exportPrevPaddingBottom = el.style.paddingBottom || '';
                el.dataset.exportPrevMarginBottom = el.style.marginBottom || '';
                el.dataset.exportPrevBoxSizing = el.style.boxSizing || '';

                el.style.boxSizing = 'border-box';
                el.style.paddingBottom = `${basePaddingBottom + extraBottomPx}px`;

                const pos = (style.position || '').toLowerCase();
                if (pos !== 'absolute' && pos !== 'fixed') {
                    el.style.marginBottom = `${baseMarginBottom - extraBottomPx}px`;
                }

                el.setAttribute(flag, '1');
                return;
            }

            if (!el.hasAttribute(flag)) {
                return;
            }

            el.style.paddingBottom = el.dataset.exportPrevPaddingBottom || '';
            el.style.marginBottom = el.dataset.exportPrevMarginBottom || '';
            el.style.boxSizing = el.dataset.exportPrevBoxSizing || '';

            delete el.dataset.exportPrevPaddingBottom;
            delete el.dataset.exportPrevMarginBottom;
            delete el.dataset.exportPrevBoxSizing;
            el.removeAttribute(flag);
        }
        """,
        {
            "targetId": target_id,
            "extraBottomPx": float(max(0.0, extra_bottom_px)),
            "enabled": bool(enabled),
        },
    )


async def collect_capture_targets(slide):
    return await slide.evaluate(
        """
        (node) => {
            const rootRect = node.getBoundingClientRect();
            const targets = [];
            let nextId = 1;
            const textTags = new Set(['H1', 'H2', 'H3', 'H4', 'H5', 'H6', 'P', 'A', 'LI', 'SPAN', 'STRONG', 'EM', 'SMALL', 'LABEL']);

            const isVisible = (el) => {
                const style = window.getComputedStyle(el);
                return style.display !== 'none' && style.visibility !== 'hidden' && Number(style.opacity || 1) > 0;
            };

            const hasTextContent = (el) => {
                const text = (el.innerText || '').replace(/\\s+/g, ' ').trim();
                return text.length > 0;
            };

            const hasVisualMedia = (el) => {
                return !!el.querySelector('img, svg, canvas, video, iframe');
            };

            const alphaFromColor = (color) => {
                if (!color) return 0;
                if (color.startsWith('rgba(')) {
                    const parts = color.slice(5, -1).split(',').map((p) => p.trim());
                    if (parts.length === 4) {
                        const a = Number(parts[3]);
                        return Number.isFinite(a) ? a : 1;
                    }
                    return 1;
                }
                if (color.startsWith('rgb(')) {
                    return 1;
                }
                if (color === 'transparent') {
                    return 0;
                }
                return 1;
            };

            const hasVisualBox = (el) => {
                const style = window.getComputedStyle(el);
                const borderW =
                    (parseFloat(style.borderTopWidth || '0') || 0) +
                    (parseFloat(style.borderRightWidth || '0') || 0) +
                    (parseFloat(style.borderBottomWidth || '0') || 0) +
                    (parseFloat(style.borderLeftWidth || '0') || 0);

                const hasBgImage = style.backgroundImage && style.backgroundImage !== 'none';
                const bgAlpha = alphaFromColor(style.backgroundColor || '');
                const hasShadow = style.boxShadow && style.boxShadow !== 'none';

                return hasBgImage || bgAlpha > 0.01 || borderW > 0.01 || !!hasShadow;
            };

            const isTransparentLayoutContainer = (el) => {
                const style = window.getComputedStyle(el);
                const display = style.display || '';
                const isLayout = display.includes('grid') || display.includes('flex') || el.tagName === 'UL' || el.tagName === 'OL';
                return isLayout && !hasVisualBox(el);
            };

            const mark = (el) => {
                const rect = el.getBoundingClientRect();
                if (rect.width < 2 || rect.height < 2) {
                    return;
                }

                const style = window.getComputedStyle(el);
                const fontSizePx = parseFloat(style.fontSize || '0') || 0;
                const tag = (el.tagName || '').toUpperCase();
                const textLike = hasTextContent(el) && (textTags.has(tag) || (!hasVisualMedia(el) && !hasVisualBox(el)));

                const desiredExtraBottom = textLike
                    ? Math.max(3, Math.min(18, fontSizePx * 0.34))
                    : 0;

                const relY = rect.top - rootRect.top;
                const roomBelow = Math.max(0, rootRect.height - (relY + rect.height));
                const extraBottomPx = Math.min(desiredExtraBottom, roomBelow);

                const id = String(nextId++);
                el.setAttribute('data-export-target-id', id);
                targets.push({
                    id,
                    x: rect.left - rootRect.left,
                    y: relY,
                    width: rect.width,
                    height: rect.height,
                    extraBottomPx,
                });
            };

            const walk = (el) => {
                if (!isVisible(el)) {
                    return;
                }

                const kids = Array.from(el.children).filter(isVisible);
                if (!kids.length) {
                    mark(el);
                    return;
                }

                if (isTransparentLayoutContainer(el)) {
                    for (const kid of kids) {
                        walk(kid);
                    }
                    return;
                }

                mark(el);
            };

            for (const child of Array.from(node.children)) {
                walk(child);
            }

            return targets;
        }
        """
    )


async def capture_slides(page, selector: str, temp_dir: pathlib.Path, include_indices: Optional[Set[int]] = None) -> List[SlideCapture]:
    locator = page.locator(selector)
    total = await locator.count()
    if total == 0:
        raise RuntimeError(f"Selector '{selector}' did not match any elements.")

    captures: List[SlideCapture] = []

    for idx in range(total):
        slide_index = idx + 1
        if include_indices and slide_index not in include_indices:
            continue

        slide = locator.nth(idx)
        await slide.scroll_into_view_if_needed()
        await page.wait_for_timeout(120)

        slide_box = await slide.bounding_box()
        if slide_box is None:
            raise RuntimeError(f"Could not resolve bounding box for slide #{slide_index}.")

        width = float(slide_box["width"])
        height = float(slide_box["height"])
        if width <= 1 or height <= 1:
            raise RuntimeError(f"Invalid slide size for #{slide_index}: {width}x{height}")

        title = await extract_slide_title(slide)
        if not title:
            title = f"Slide {slide_index:02d}"

        layers: List[LayerCapture] = []

        await hide_direct_children(slide)
        bg_path = temp_dir / f"slide_{slide_index:02d}_layer_000_background.png"
        await slide.screenshot(path=str(bg_path), type="png", animations="disabled")
        await restore_direct_children(slide)

        layers.append(
            LayerCapture(
                image_path=bg_path,
                x=0.0,
                y=0.0,
                width=width,
                height=height,
            )
        )

        await clear_capture_targets(slide)
        targets = await collect_capture_targets(slide)

        for target_idx, target in enumerate(targets, start=1):
            layer_path = temp_dir / f"slide_{slide_index:02d}_layer_{target_idx:03d}.png"
            target_id = str(target.get("id", "")).strip()
            if not target_id:
                continue

            rel_x = float(target.get("x", 0.0))
            rel_y = float(target.get("y", 0.0))
            target_width = float(target.get("width", 0.0))
            target_height = float(target.get("height", 0.0))
            extra_bottom = float(target.get("extraBottomPx", 0.0))

            if target_width <= 1 or target_height <= 1:
                continue

            target_locator = slide.locator(f'[data-export-target-id="{target_id}"]')

            try:
                if extra_bottom > 0:
                    await set_target_extra_bottom_padding(slide, target_id, extra_bottom, True)
                await target_locator.screenshot(path=str(layer_path), type="png", animations="disabled")
            except Exception:
                continue
            finally:
                if extra_bottom > 0:
                    await set_target_extra_bottom_padding(slide, target_id, extra_bottom, False)

            layers.append(
                LayerCapture(
                    image_path=layer_path,
                    x=rel_x,
                    y=rel_y,
                    width=target_width,
                    height=target_height + extra_bottom,
                )
            )

        await clear_capture_targets(slide)

        captures.append(
            SlideCapture(
                source_index=slide_index,
                title=title,
                width=width,
                height=height,
                layers=layers,
            )
        )

        print(
            f"Captured slide {slide_index}/{total}: {int(width)}x{int(height)}, "
            f"layers={len(layers)}, title='{title}'"
        )

    return captures


def build_pptx(captures: List[SlideCapture], output_path: pathlib.Path):
    if not captures:
        raise RuntimeError("No slides were captured.")

    first = captures[0]
    ratio = first.width / first.height
    target_height_in = 7.5
    target_width_in = target_height_in * ratio

    presentation = Presentation()
    presentation.slide_width = Inches(target_width_in)
    presentation.slide_height = Inches(target_height_in)

    blank_layout = presentation.slide_layouts[6]

    for capture in captures:
        slide = presentation.slides.add_slide(blank_layout)

        scale_x = target_width_in / capture.width
        scale_y = target_height_in / capture.height

        for layer in capture.layers:
            left = Inches(layer.x * scale_x)
            top = Inches(layer.y * scale_y)
            width = Inches(layer.width * scale_x)
            height = Inches(layer.height * scale_y)
            slide.shapes.add_picture(str(layer.image_path), left, top, width=width, height=height)

    if output_path.exists():
        output_path.unlink()

    presentation.save(str(output_path))


async def export_canva_pptx(
    client_slug: str,
    project_slug: str,
    selector: Optional[str],
    output_name: Optional[str],
    port: int,
    wait_ms: int,
    slide_indices: Optional[Set[int]],
):
    if not project_exists(client_slug, project_slug):
        raise FileNotFoundError(
            f"Project not found: public/clients/{client_slug}/projects/{project_slug}/index.html"
        )

    os.chdir(PUBLIC)

    if TMP_ROOT.exists():
        shutil.rmtree(TMP_ROOT)
    TMP_ROOT.mkdir(parents=True, exist_ok=True)

    server = start_server(port)
    time.sleep(0.3)

    url = f"http://127.0.0.1:{port}/clients/{client_slug}/projects/{project_slug}/index.html"

    try:
        async with async_playwright() as playwright:
            browser = await playwright.chromium.launch()
            page = await browser.new_page(
                viewport={"width": 2600, "height": 1600},
                device_scale_factor=2,
            )

            print(f"Loading: {url}")
            await page.goto(url, wait_until="networkidle")
            await wait_for_assets(page)
            await page.wait_for_timeout(wait_ms)

            chosen_selector = selector or await detect_default_selector(page)
            if not chosen_selector:
                raise RuntimeError("Could not auto-detect slide selector. Pass one with --selector.")

            page_title = (await page.title()).strip()
            output_filename = sanitize_filename(output_name or f"{page_title} - Canva Editable", ".pptx")
            output_path = ROOT / output_filename

            print(f"Using selector: {chosen_selector}")
            captures = await capture_slides(page, chosen_selector, TMP_ROOT, include_indices=slide_indices)

            if not captures:
                if slide_indices:
                    selection = ", ".join(str(idx) for idx in sorted(slide_indices))
                    raise RuntimeError(f"No slides matched --slide-indices {selection}")
                raise RuntimeError("No slides were captured.")

            build_pptx(captures, output_path)
            await browser.close()

            print(f"PPTX exported: {output_path}")
            print(f"Slides exported: {len(captures)}")
            return output_path
    finally:
        server.shutdown()
        server.server_close()


def parse_slide_indices(raw: Optional[str]) -> Optional[Set[int]]:
    if not raw:
        return None

    result: Set[int] = set()
    for chunk in str(raw).split(','):
        token = chunk.strip()
        if not token:
            continue

        if '-' in token:
            start_raw, end_raw = token.split('-', 1)
            start = int(start_raw.strip())
            end = int(end_raw.strip())
            if start <= 0 or end <= 0:
                raise ValueError("Slide indices must be positive integers.")
            lo, hi = (start, end) if start <= end else (end, start)
            for value in range(lo, hi + 1):
                result.add(value)
            continue

        value = int(token)
        if value <= 0:
            raise ValueError("Slide indices must be positive integers.")
        result.add(value)

    return result or None


def parse_args():
    parser = argparse.ArgumentParser(description="Export static project slides to Canva-compatible PPTX")
    parser.add_argument("--client", default="cbhn", help="Client slug under public/clients (default: cbhn)")
    parser.add_argument("--project", required=True, help="Project slug under public/clients/<client>/projects")
    parser.add_argument("--selector", help="CSS selector for slide containers (e.g. .slide)")
    parser.add_argument("--output", help="Optional output PPTX filename")
    parser.add_argument("--slide-indices", help="Optional slide numbers (e.g. 1,2,3,4 or 1-4)")
    parser.add_argument("--port", type=int, default=8001, help="Local preview port")
    parser.add_argument("--wait-ms", type=int, default=1200, help="Extra wait before capture (ms)")
    return parser.parse_args()


if __name__ == "__main__":
    arguments = parse_args()
    indices = parse_slide_indices(arguments.slide_indices)
    asyncio.run(
        export_canva_pptx(
            arguments.client,
            arguments.project,
            arguments.selector,
            arguments.output,
            arguments.port,
            arguments.wait_ms,
            indices,
        )
    )
